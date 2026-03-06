"""PPTX renderer — generates PowerPoint presentations using python-pptx."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from design_system.models import (
    DesignTokens,
    Presentation,
    RenderResult,
    SlideContent,
    SlideType,
    ThemeInfo,
)
from design_system.renderers.base import BaseRenderer

logger = logging.getLogger(__name__)

# Slide dimensions: Widescreen 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout margins
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
CONTENT_WIDTH = Inches(11.733)
CONTENT_HEIGHT = Inches(5.5)

# Footer / branding bar
FOOTER_HEIGHT = Inches(0.35)
FOOTER_TOP = SLIDE_HEIGHT - FOOTER_HEIGHT


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lighten(hex_color: str, factor: float = 0.92) -> RGBColor:
    """Lighten a hex color toward white by the given factor (0=original, 1=white)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return RGBColor(r, g, b)


class PptxRenderer(BaseRenderer):
    """Render presentations as PPTX files using python-pptx."""

    @property
    def format_name(self) -> str:
        return "pptx"

    def render(
        self,
        presentation: Presentation,
        theme: ThemeInfo,
        output_path: Path,
    ) -> RenderResult:
        """Render a Presentation to a .pptx file."""
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        tokens = theme.tokens
        total = len(presentation.slides)

        for idx, slide_content in enumerate(presentation.slides):
            self._add_slide(prs, slide_content, tokens, idx + 1, total)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("Saved PPTX: %s (%d slides)", output_path, len(presentation.slides))

        return RenderResult(
            output_path=output_path,
            format="pptx",
            theme=theme.name,
            slide_count=len(presentation.slides),
        )

    def _add_slide(
        self,
        prs: PptxPresentation,
        content: SlideContent,
        tokens: DesignTokens,
        slide_num: int,
        total: int,
    ) -> None:
        """Add a single slide based on its type."""
        layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(layout)

        handlers = {
            SlideType.COVER: self._render_cover,
            SlideType.SECTION: self._render_section,
            SlideType.BULLET: self._render_bullet,
            SlideType.TABLE: self._render_table,
            SlideType.QUOTE: self._render_quote,
            SlideType.TEXT: self._render_text,
            SlideType.CLOSING: self._render_closing,
        }

        handler = handlers.get(content.slide_type, self._render_text)
        handler(slide, content, tokens)

        # Add footer/branding to content slides (not cover/closing)
        if content.slide_type not in (SlideType.COVER, SlideType.CLOSING):
            self._add_footer(slide, tokens, slide_num, total)

    # ----- Slide type renderers -----

    def _render_cover(self, slide, content: SlideContent, tokens: DesignTokens) -> None:
        """Cover slide with full-width colored background and accent stripe."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)

        # Top-left decorative bar (secondary color)
        deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            SLIDE_WIDTH,
            Inches(0.15),
        )
        deco_fill = deco.fill
        deco_fill.solid()
        deco_fill.fore_color.rgb = _hex_to_rgb(tokens.color.secondary)
        deco.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(2.0),
            CONTENT_WIDTH,
            Inches(1.8),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(tokens.typography.title_size + 12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = tokens.typography.font_family
        p.alignment = PP_ALIGN.LEFT

        # Accent line between title and subtitle
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_LEFT,
            Inches(3.95),
            Inches(3),
            Pt(4),
        )
        accent_fill = accent.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        accent.line.fill.background()

        # Subtitle
        if content.subtitle:
            sub_box = slide.shapes.add_textbox(
                MARGIN_LEFT,
                Inches(4.3),
                CONTENT_WIDTH,
                Inches(1.0),
            )
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = content.subtitle
            p2.font.size = Pt(tokens.typography.subheading_size + 2)
            p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p2.font.name = tokens.typography.font_family

        # Bottom bar
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(7.0),
            SLIDE_WIDTH,
            Inches(0.5),
        )
        bottom_fill = bottom.fill
        bottom_fill.solid()
        bottom_fill.fore_color.rgb = _hex_to_rgb(tokens.color.secondary)
        bottom.line.fill.background()

    def _render_section(
        self, slide, content: SlideContent, tokens: DesignTokens
    ) -> None:
        """Section divider slide — bold accent stripe at top."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(tokens.color.surface)

        # Full-width accent bar at top
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            SLIDE_WIDTH,
            Inches(0.12),
        )
        bar_fill = bar.fill
        bar_fill.solid()
        bar_fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)
        bar.line.fill.background()

        # Left vertical accent bar
        vbar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(2.5),
            Pt(6),
            Inches(2.5),
        )
        vbar_fill = vbar.fill
        vbar_fill.solid()
        vbar_fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)
        vbar.line.fill.background()

        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1.2),
            Inches(2.5),
            Inches(10),
            Inches(2.5),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(tokens.typography.heading_size + 8)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(tokens.color.text)
        p.font.name = tokens.typography.font_family

    def _render_bullet(
        self, slide, content: SlideContent, tokens: DesignTokens
    ) -> None:
        """Bullet list slide — single text frame, no overlapping shapes."""
        self._add_slide_background(slide, tokens)
        self._add_title_bar(slide, content.title, tokens)

        body_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.05),
            Inches(1.8),
            CONTENT_WIDTH - Inches(0.05),
            Inches(5.0),
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        primary_rgb = _hex_to_rgb(tokens.color.primary)
        text_rgb = _hex_to_rgb(tokens.color.text)
        bullet_size = Pt(tokens.typography.body_size)

        for i, bullet in enumerate(content.bullets or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Colored bullet marker as inline text
            run_marker = p.add_run()
            run_marker.text = "\u25A0  "  # filled square
            run_marker.font.size = Pt(tokens.typography.body_size - 4)
            run_marker.font.color.rgb = primary_rgb
            run_marker.font.name = tokens.typography.font_family
            # Bullet text
            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = bullet_size
            run_text.font.color.rgb = text_rgb
            run_text.font.name = tokens.typography.font_family
            p.space_after = Pt(tokens.spacing.sm + 4)

    def _render_table(self, slide, content: SlideContent, tokens: DesignTokens) -> None:
        """Table slide."""
        self._add_slide_background(slide, tokens)
        self._add_title_bar(slide, content.title, tokens)

        if not content.table:
            return

        rows = len(content.table.rows) + 1  # +1 for header
        cols = len(content.table.headers)
        if cols == 0:
            return

        table_width = min(CONTENT_WIDTH, Inches(12))
        table_height = Inches(min(0.45 * rows, 5.0))

        table_shape = slide.shapes.add_table(
            rows,
            cols,
            MARGIN_LEFT,
            Inches(1.8),
            table_width,
            table_height,
        )
        table = table_shape.table

        # Style header row
        for j, header in enumerate(content.table.headers):
            cell = table.cell(0, j)
            cell.text = header
            self._style_cell(cell, tokens, is_header=True)

        # Style data rows
        for i, row in enumerate(content.table.rows):
            for j, value in enumerate(row):
                if j < cols:
                    cell = table.cell(i + 1, j)
                    cell.text = value
                    self._style_cell(cell, tokens, is_header=False, row_idx=i)

    def _render_quote(self, slide, content: SlideContent, tokens: DesignTokens) -> None:
        """Quote/key message slide with accent bar."""
        self._add_slide_background(slide, tokens)

        # Left accent bar (thick, primary color)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0),
            Inches(1.8),
            Pt(6),
            Inches(3.5),
        )
        accent_fill = accent.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)
        accent.line.fill.background()

        # Large quote mark
        quote_mark = slide.shapes.add_textbox(
            Inches(1.4),
            Inches(1.2),
            Inches(2),
            Inches(1.5),
        )
        tf_q = quote_mark.text_frame
        p_q = tf_q.paragraphs[0]
        p_q.text = "\u201c"
        p_q.font.size = Pt(96)
        p_q.font.color.rgb = _hex_to_rgb(tokens.color.primary)
        p_q.font.name = tokens.typography.font_family

        # Quote text
        quote_box = slide.shapes.add_textbox(
            Inches(1.8),
            Inches(2.5),
            Inches(9.5),
            Inches(3),
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.quote
        p.font.size = Pt(tokens.typography.subheading_size + 2)
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb(tokens.color.text)
        p.font.name = tokens.typography.font_family
        p.alignment = PP_ALIGN.LEFT

    def _render_text(self, slide, content: SlideContent, tokens: DesignTokens) -> None:
        """General text/body slide."""
        self._add_slide_background(slide, tokens)
        self._add_title_bar(slide, content.title, tokens)

        body_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(1.8),
            CONTENT_WIDTH,
            CONTENT_HEIGHT,
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.body
        p.font.size = Pt(tokens.typography.body_size)
        p.font.color.rgb = _hex_to_rgb(tokens.color.text)
        p.font.name = tokens.typography.font_family
        p.space_after = Pt(tokens.spacing.md)

    def _render_closing(
        self, slide, content: SlideContent, tokens: DesignTokens
    ) -> None:
        """Closing/thank you slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)

        # Top decorative bar (secondary)
        deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            SLIDE_WIDTH,
            Inches(0.15),
        )
        deco_fill = deco.fill
        deco_fill.solid()
        deco_fill.fore_color.rgb = _hex_to_rgb(tokens.color.secondary)
        deco.line.fill.background()

        title_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(2.8),
            CONTENT_WIDTH,
            Inches(2),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(tokens.typography.title_size + 8)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = tokens.typography.font_family
        p.alignment = PP_ALIGN.CENTER

        # Subtitle (e.g. branding line)
        if content.subtitle:
            sub_box = slide.shapes.add_textbox(
                MARGIN_LEFT,
                Inches(4.5),
                CONTENT_WIDTH,
                Inches(0.8),
            )
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = content.subtitle
            p2.font.size = Pt(tokens.typography.subheading_size)
            p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p2.font.name = tokens.typography.font_family
            p2.alignment = PP_ALIGN.CENTER

        # Bottom bar
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(7.0),
            SLIDE_WIDTH,
            Inches(0.5),
        )
        bottom_fill = bottom.fill
        bottom_fill.solid()
        bottom_fill.fore_color.rgb = _hex_to_rgb(tokens.color.secondary)
        bottom.line.fill.background()

    # ----- Helper methods -----

    def _add_slide_background(self, slide, tokens: DesignTokens) -> None:
        """Set slide background color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(tokens.color.background)

    def _add_title_bar(self, slide, title: str, tokens: DesignTokens) -> None:
        """Add a title bar at the top of the slide with accent line."""
        if not title:
            return

        # Top-edge accent stripe (full width, thin)
        top_stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            SLIDE_WIDTH,
            Pt(4),
        )
        top_fill = top_stripe.fill
        top_fill.solid()
        top_fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)
        top_stripe.line.fill.background()

        # Accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_LEFT,
            Inches(1.35),
            Inches(2.0),
            Pt(3),
        )
        line_fill = line.fill
        line_fill.solid()
        line_fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)
        line.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(0.4),
            CONTENT_WIDTH,
            Inches(0.85),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(tokens.typography.heading_size)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(tokens.color.text)
        p.font.name = tokens.typography.font_family

    def _add_footer(
        self, slide, tokens: DesignTokens, slide_num: int, total: int
    ) -> None:
        """Add branding footer bar with slide number."""
        # Footer background bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            FOOTER_TOP,
            SLIDE_WIDTH,
            FOOTER_HEIGHT,
        )
        bar_fill = bar.fill
        bar_fill.solid()
        bar_fill.fore_color.rgb = _hex_to_rgb(tokens.color.secondary)
        bar.line.fill.background()

        # Slide number (right-aligned)
        num_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(1.5),
            FOOTER_TOP + Pt(2),
            Inches(1.2),
            FOOTER_HEIGHT - Pt(4),
        )
        tf = num_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{slide_num} / {total}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = tokens.typography.font_family
        p.alignment = PP_ALIGN.RIGHT

    def _style_cell(
        self,
        cell,
        tokens: DesignTokens,
        *,
        is_header: bool = False,
        row_idx: int = 0,
    ) -> None:
        """Style a table cell."""
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(tokens.typography.caption_size + 1)
            paragraph.font.name = tokens.typography.font_family

            if is_header:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                paragraph.font.color.rgb = _hex_to_rgb(tokens.color.text)

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if is_header:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(tokens.color.primary)
        elif row_idx % 2 == 1:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(tokens.color.surface)
